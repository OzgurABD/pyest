import os
from io import BytesIO
from fastapi import UploadFile
import fitz  # PyMuPDF
from pdf2docx import Converter  # pdf2docx
from docx import Document
from pptx import Presentation
from pptx.util import Inches

TEMP_DIR = "temp_files"

if not os.path.exists(TEMP_DIR):
    os.makedirs(TEMP_DIR)


def extract_text_from_pdf(pdf_file: BytesIO) -> str:
    doc = fitz.open(pdf_file)
    text = ""
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        text += page.get_text()
    return text


def create_docx_from_text(text: str, output_path: str) -> None:
    doc = Document()
    doc.add_paragraph(text)
    doc.save(output_path)


async def handle_pdf_word_v1(file: UploadFile) -> any:
    pdf_path = os.path.join(TEMP_DIR, file.filename)
    with open(pdf_path, "wb") as f:
        f.write(await file.read())

    text = extract_text_from_pdf(pdf_path)

    docx_filename = f"{file.filename.split('.')[0]}.docx"
    docx_path = os.path.join(TEMP_DIR, docx_filename)
    create_docx_from_text(text, docx_path)

    return {"fileName": docx_filename, "filePath": docx_path}


async def handle_pdf_word_v2(file: UploadFile) -> any:
    pdf_path = os.path.join(TEMP_DIR, file.filename)
    with open(pdf_path, "wb") as f:
        f.write(await file.read())

    docx_filename = f"{file.filename.split('.')[0]}.docx"
    docx_path = os.path.join(TEMP_DIR, docx_filename)

    converter = Converter(pdf_path)
    converter.convert(docx_path, start=0, end=None)

    # os.remove(pdf_path)

    return {"fileName": docx_filename, "filePath": docx_path}


def extract_text_and_images_from_pdf(pdf_file: BytesIO):
    doc = fitz.open(pdf_file)
    text = ""
    images = []
    for page_num in range(doc.page_count):
        page = doc.load_page(page_num)
        text += page.get_text()
        for img_index, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            images.append(image_bytes)
    return text, images


async def handle_pdf_to_pptx(file: UploadFile) -> str:
    pdf_path = os.path.join(TEMP_DIR, file.filename)
    with open(pdf_path, "wb") as f:
        f.write(await file.read())

    text, images = extract_text_and_images_from_pdf(pdf_path)

    pptx_filename = f"{file.filename.split('.')[0]}.pptx"
    pptx_path = os.path.join(TEMP_DIR, pptx_filename)

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    content = slide.shapes.placeholders[1]
    title.text = "PDF Dönüştürme"
    content.text = text

    for img_index, image_bytes in enumerate(images):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.add_picture(
            BytesIO(image_bytes), Inches(1), Inches(1), width=Inches(8), height=Inches(6))

    prs.save(pptx_path)

    return pptx_path
